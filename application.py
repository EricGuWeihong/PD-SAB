import streamlit as st
from streamlit_text_rating.st_text_rater import st_text_rater
from streamlit.components.v1 import html
import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document
import io
import broadscope_bailian

import dashscope
from dashscope import Files
from dashscope import Assistants
from dashscope import Application
from dashscope.assistants.files import Files as AssistantFiles
from dashscope import Threads
from dashscope import Messages
from dashscope import Runs
import os
import json
import time 

# Configure credentials
accessKeyId = st.secrets["ALIYUN_ACCESS_KEY"]
accessKeySecret = st.secrets["ALIYUN_ACCESS_Key_SECRET"]
agentKey = st.secrets["ALIYUN_AGENT_KEY"]
dashscope.api_key=st.secrets["DASHSCOPE_API_KEY"]

st.title("💬 Sales AI Buddy")
st.caption("🚀 派得 - 智能销售伙伴")

with st.sidebar:
    with st.expander("参数设置"):
        model = st.radio(label="选择模型",options=["智能陪练", "企业知识库","通义千问"],help="Agent为智能体；Max更智能，但费用较高")
        if model == "智能陪练":
            app_id = st.secrets["ALIYUN_APP_ID_APP"]
        elif model == "企业知识库":
            app_id = st.secrets["ALIYUN_APP_ID_MAX"]
        else:
            app_id = st.secrets["ALIYUN_APP_ID"]

        temperature = st.slider(label="温度",min_value=0.01,max_value=1.0, step=0.1, help="数值越大，创造力越大但相应的可能会编撰")
    labels = st.multiselect(label="知识标签",options=["产品","准入"])
    upload_file = st.file_uploader("上传文档", type=("txt", "pdf", "ppt","pptx","doc","docx"))

def clear_chat_history():
    if model == "智能陪练":
        st.session_state["messages"] = [{"role":"assistant","content":"请选择一个场景，回复数字开始模拟对话练习。对话过程中，回复【结束】中止模拟练习。\n1. BRCA阴性且HRD阳性; \n2. 联合贝伐珠单抗治疗; \n3. 没有使用贝伐珠单抗; \n4. 随机场景"}]
    elif model == "企业知识库":
        st.session_state["messages"] = [{"role":"assistant","content":"有什么问题想问我？"}]
    else:
        st.session_state["messages"] = [{"role":"assistant","content":"嘿，我有什么可以帮忙的？"}]
    st.session_state["thread"] = None
    st.session_state["label"] = True

st.sidebar.button(label="清空聊天记录",on_click=clear_chat_history)

def get_app_response(prompt):
    if st.session_state["thread"]:
        resps = Application.call(app_id=app_id,
                                prompt=prompt,
                                stream=True,
                                session_id = st.session_state["thread"]
                                )
    else:
        resps = Application.call(app_id=app_id,
                                prompt=prompt,
                                stream=True
                                )
    return resps

if "messages" not in st.session_state:
    clear_chat_history()

for msg in st.session_state.messages:
    st.chat_message(msg["role"]).markdown(msg["content"])

if prompt := st._bottom.chat_input("请输入..."):
    st.session_state.messages.append({"role": "user", "content": prompt})
    with st.chat_message("user"):
        st.markdown(prompt)

    if upload_file:
        if upload_file.type == "text/plain":
            # 处理文本文件
            article = upload_file.read().decode('utf-8')
        elif upload_file.type == "application/pdf":
            # 处理PDF文件
            with fitz.open(stream=upload_file.read(), filetype="pdf") as doc:
                text = ""
                for page in doc:
                    text += page.get_text()
                article = text
        elif upload_file.type in ["application/vnd.ms-powerpoint", "application/vnd.openxmlformats-officedocument.presentationml.presentation"]:
            # 处理PPT和PPTX文件
            ppt = Presentation(io.BytesIO(upload_file.read()))
            text = ""
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            article = text
        elif upload_file.type in ["application/msword", "application/vnd.openxmlformats-officedocument.wordprocessingml.document"]:
            # 处理DOC和DOCX文件
            doc = Document(io.BytesIO(upload_file.read()))
            text = ""
            for para in doc.paragraphs:
                text += para.text + "\n"
            article = text
        # prompt = prompt + "\n以下是用户上传的文件或文件:\n" + upload_file.read().decode()
        prompt = prompt + "\n以下是用户上传的文件:\n" + article

    with st.chat_message("assistant"):
        if model == "智能陪练":
            resps = get_app_response(prompt)

            with st.spinner("CPU飞速运转中..."):
                with st.empty():
                    new_text = ""
                    for resp in resps:
                        if resp.status_code == 200:
                            new_text = resp.output.text
                            st.markdown(new_text)
                        else:
                            st.warning(resp.output)
                    st.session_state["thread"] = resp.output.session_id
        else:
            client = broadscope_bailian.AccessTokenClient(accessKeyId, accessKeySecret, agentKey);
            token = client.get_token();
            with st.spinner("CPU飞速运转中..."):
                resp = broadscope_bailian.Completions(token=token).create(
                        app_id=app_id,
                        prompt=prompt,
                        history=st.session_state.messages,
                        # 设置模型参数temperature
                        temperature=temperature,
                        # 流式输出，默认False
                        stream=True,
                        # 开启增量输出模式，后面输出不会包含已经输出的内容
                        # incremental_output=True,
                        # 返回choice message结果
                        result_format="message",
                        # 返回文档检索的文档引用数据, 传入为simple或indexed
                        doc_reference_type="indexed",
                        # 文档标签code列表
                        doc_tag_codes=[],
                        # 是否输出文档召回推理过程相关信息
                        has_thoughts = True,
                        # 30秒超时
                        timeout=30,
                    )
            
            if not next(resp)["Success"]:
                st.warning("对不起😵‍💫, 请求ID: %s, 错误码: %s, 错误信息: %s" %
                    (resp.get("RequestId"), resp.get("Code"), resp.get("Message")))

            with st.empty():
                new_text = ""
                for chunk in resp:
                    if "Text" in chunk["Data"]:
                        new_text = chunk["Data"]["Text"]
                        st.markdown(new_text,unsafe_allow_html=True)

            resp_data = chunk
            doc_refs = chunk["Data"]["DocReferences"]
            if doc_refs is not None and len(doc_refs) > 0:
                for i in range(0, len(doc_refs)):
                    doc_name = doc_refs[i]["IndexId"] + " - " + doc_refs[i]["DocName"]
                    st.page_link(label=doc_name,page=doc_refs[i]["DocUrl"],icon="📑",help=doc_refs[i]["Text"])
                # st.info("参考文档：%s" % doc_refs[0]["DocName"])
                # st.markdown(doc_refs, unsafe_allow_html=True)
            st_text_rater(text="我的回复有帮助到你吗?",color_text="rgb(100,100,100)",font_size="12px")
    st.session_state.messages.append({"role": "assistant", "content": new_text})
    # with st.sidebar.expander(label="输出数据",expanded=False):
    #     st.write(resp_data)
    #     st.write(st.session_state.messages)