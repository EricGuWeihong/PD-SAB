import streamlit as st
from streamlit_text_rating.st_text_rater import st_text_rater
from streamlit.components.v1 import html
import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document
import io
import broadscope_bailian

import dashscope
from dashscope import Files
from dashscope import Assistants
from dashscope.assistants.files import Files as AssistantFiles
from dashscope import Threads
from dashscope import Messages
from dashscope import Runs
import os
import json
import time 

# Configure credentials
accessKeyId = st.secrets["ALIYUN_ACCESS_KEY"]
accessKeySecret = st.secrets["ALIYUN_ACCESS_Key_SECRET"]
agentKey = st.secrets["ALIYUN_AGENT_KEY"]
dashscope.api_key=st.secrets["DASHSCOPE_API_KEY"]
endpoint = 'oss-cn-beijing.aliyuncs.com'
# bucket_name = 'eric-paide-kb'

# Retrieve the existing assistant
assistant_id = st.secrets["ALIYUN_APP_ID_AGENT"]
assistant = Assistants.retrieve(assistant_id)

# #==== Start of Upload KB files to AssistantFile, only need to run once; If any updates, please re-run this code ====
# # Get all files in local KB folder
# local_folder_path = 'KB'  # Replace with the actual local folder path
# file_list = os.listdir(local_folder_path)

# # Upload files to AssistantFiles and get file IDs
# file_ids = []
# for file_name in file_list:
#     local_file_path = os.path.join(local_folder_path, file_name)
#     file_id = Files.upload(file_path=local_file_path, purpose='assistant', description='Assistant KB')['output']['uploaded_files'][0]['file_id']
#     file_ids.append(file_id)

# # Create AssistantFiles object based on the uploaded oss files
# for file_id in file_ids:
#     assistant_file = AssistantFiles.create(assistant_id, file_id=file_id)
#     print("status code:", assistant_file.status_code, " - id:", assistant_file.id)
# #==== End of uploading KB files ====

# List all the files in the assistant
assistant_files = AssistantFiles.list(assistant_id,
                                      sort_by="created_at",
                                      sort_order="desc")
print(len(assistant_files.data))
print(assistant_files.data)

# #==== Start of Delete all the files in the assistant
# for assistant_file in assistant_files.data:
#     assistant_file_deleted = AssistantFiles.delete(assistant_file.id, assistant_id=assistant_id)
#     print(assistant_file_deleted)
# #==== End of Delete all the files in the assistant

def clear_chat_history():
    st.session_state["messages"] = [{"role":"assistant","content":"嘿，我有什么可以帮忙的？"}]
    st.session_state["thread"] = None
st.sidebar.button(label="清空聊天记录",on_click=clear_chat_history)

if "messages" not in st.session_state:
    clear_chat_history()

with st.sidebar:
    with st.expander("参数设置"):
        model = st.radio(label="选择模型",options=["智能陪练", "企业知识库","通义千问"],help="Agent为智能体；Max更智能，但费用较高")
        if model == "智能陪练":
            app_id = st.secrets["ALIYUN_APP_ID_AGENT"]
        elif model == "企业知识库":
            app_id = st.secrets["ALIYUN_APP_ID_MAX"]
        else:
            app_id = st.secrets["ALIYUN_APP_ID"]

        temperature = st.slider(label="温度",min_value=0.01,max_value=1.0, step=0.1, help="数值越大，创造力越大但相应的可能会编撰")
    labels = st.multiselect(label="知识标签",options=["产品","准入"])
    upload_file = st.file_uploader("上传文档", type=("txt", "pdf", "ppt","pptx","doc","docx"))

st.title("💬 Sales AI Buddy")
st.caption("🚀 派得 - 智能销售伙伴")


for msg in st.session_state.messages:
    st.chat_message(msg["role"]).write(msg["content"])

st._bottom.info("AI可能会犯错误，请核实重要信息",icon="🚨")
if prompt := st._bottom.chat_input("请输入你的问题"):
    st.session_state.messages.append({"role": "user", "content": prompt})
    with st.chat_message("user"):
        st.markdown(prompt)

    if upload_file:
        if upload_file.type == "text/plain":
            # 处理文本文件
            article = upload_file.read().decode('utf-8')
        elif upload_file.type == "application/pdf":
            # 处理PDF文件
            with fitz.open(stream=upload_file.read(), filetype="pdf") as doc:
                text = ""
                for page in doc:
                    text += page.get_text()
                article = text
        elif upload_file.type in ["application/vnd.ms-powerpoint", "application/vnd.openxmlformats-officedocument.presentationml.presentation"]:
            # 处理PPT和PPTX文件
            ppt = Presentation(io.BytesIO(upload_file.read()))
            text = ""
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            article = text
        elif upload_file.type in ["application/msword", "application/vnd.openxmlformats-officedocument.wordprocessingml.document"]:
            # 处理DOC和DOCX文件
            doc = Document(io.BytesIO(upload_file.read()))
            text = ""
            for para in doc.paragraphs:
                text += para.text + "\n"
            article = text
        # prompt = prompt + "\n以下是用户上传的文件或文件:\n" + upload_file.read().decode()
        prompt = prompt + "\n以下是用户上传的文件:\n" + article

    with st.chat_message("assistant"):
        if model == "智能陪练":
            if st.session_state["thread"] == None:
                thread = Threads.create(messages=[{"role": "user", "content": prompt}])
                # Check if the thread was created successfully
                if thread.status_code != 200:
                    st.error('Create Thread failed, status code: %s, code: %s, message: %s' % (thread.status_code, thread.code, thread.message))
                else:
                    print('Create Thread success, id: %s' % thread.id)
                    st.session_state["thread"] = thread.id
            else:
                message = Messages.create(thread_id=st.session_state["thread"], role="user", content=prompt)
                if message.status_code != 200:
                    st.error('Create Message failed, status code: %s, code: %s, message: %s' % (message.status_code, message.code, message.message))
                else:
                    print('Create Message success, id: %s' % message.id)
            
            # Create a new run to run message
            run = Runs.create(thread_id=st.session_state["thread"], assistant_id=assistant_id)
            if run.status_code != 200:
                st.error('Create Run failed, status code: %s, code: %s, message: %s' % (run.status_code, run.status, run.message))
            else:
                print('Create Run success, id: %s' % run.id)

            with st.spinner("CPU飞速运转中..."):
                # Wait for the run to complete or requires_action
                msgs = []
                while True:
                    run = Runs.wait(run.id, thread_id=st.session_state["thread"])
                    if run.status_code != 200:
                        st.error('Wait for Run failed, status code: %s, code: %s' % (run.status_code, run.status))
                    else:
                        if run.status == "completed":
                            msgs = Messages.list(thread_id=st.session_state["thread"],limit=1,order="desc")
                            if msgs.status_code != 200:
                                st.error('Get run failed, status code: %s, code: %s, message: %s' % (msgs.status_code, run.status, msgs.message))
                            else:
                                print(json.dumps(msgs, default=lambda o: o.__dict__, sort_keys=True, indent=4))
                            break
                        if run.status == "failed" or run.status == "expired":
                            st.error('Run failed, status code: %s, code: %s' % (run.status_code, run.status))
                            print(run)
                            break
                        else:
                            time.sleep(1)
            
            # Get the thread messages, to get the run output
            new_text = ""
            with st.empty():
                if msgs:
                    for msg in msgs.data:
                        print(msg)
                        new_text = msg.content[0]["text"]["value"]
                        st.markdown(new_text,unsafe_allow_html=True)
            
                    # st_text_rater(text="我的回复有帮助到你吗?",color_text="rgb(100,100,100)",font_size="12px")
            
        else:
            client = broadscope_bailian.AccessTokenClient(accessKeyId, accessKeySecret, agentKey);
            token = client.get_token();
            with st.spinner("CPU飞速运转中..."):
                resp = broadscope_bailian.Completions(token=token).create(
                        app_id=app_id,
                        prompt=prompt,
                        history=st.session_state.messages,
                        # 设置模型参数temperature
                        temperature=temperature,
                        # 流式输出，默认False
                        stream=True,
                        # 开启增量输出模式，后面输出不会包含已经输出的内容
                        # incremental_output=True,
                        # 返回choice message结果
                        result_format="message",
                        # 返回文档检索的文档引用数据, 传入为simple或indexed
                        doc_reference_type="indexed",
                        # 文档标签code列表
                        doc_tag_codes=[],
                        # 是否输出文档召回推理过程相关信息
                        has_thoughts = True,
                        # 30秒超时
                        timeout=30,
                    )
            
            if not next(resp)["Success"]:
                st.warning("对不起😵‍💫, 请求ID: %s, 错误码: %s, 错误信息: %s" %
                    (resp.get("RequestId"), resp.get("Code"), resp.get("Message")))

            with st.empty():
                new_text = ""
                for chunk in resp:
                    if "Text" in chunk["Data"]:
                        new_text = chunk["Data"]["Text"]
                        st.markdown(new_text,unsafe_allow_html=True)

            resp_data = chunk
            doc_refs = chunk["Data"]["DocReferences"]
            if doc_refs is not None and len(doc_refs) > 0:
                for i in range(0, len(doc_refs)):
                    doc_name = doc_refs[i]["IndexId"] + " - " + doc_refs[i]["DocName"]
                    st.page_link(label=doc_name,page=doc_refs[i]["DocUrl"],icon="📑",help=doc_refs[i]["Text"])
                # st.info("参考文档：%s" % doc_refs[0]["DocName"])
                # st.markdown(doc_refs, unsafe_allow_html=True)
            st_text_rater(text="我的回复有帮助到你吗?",color_text="rgb(100,100,100)",font_size="12px")
    st.session_state.messages.append({"role": "assistant", "content": new_text})
    # with st.sidebar.expander(label="输出数据",expanded=False):
    #     st.write(resp_data)
    #     st.write(st.session_state.messages)